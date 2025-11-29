# core.py - ИСПРАВЛЕННАЯ ВЕРСИЯ
from .old_functions import (
    ContentElement, ContentType, ElementStyle, TextStyle, BorderStyle, 
    SizeConstraints, LayoutStrategy, ColorThemes, Inches, Pt, RGBColor,
    PP_ALIGN, MSO_SHAPE
)
from .layout_engine import SmartLayoutEngine
from .text_module import AdvancedTextModule
from .graphics_module import GraphicsBuilder
from .media_module import MediaManager
from .styles_module import StyleSystem
from pptx import Presentation
from pptx.enum.text import PP_ALIGN as PPTX_PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE as PPTX_MSO_SHAPE
from typing import List, Dict, Optional, Tuple, Union

class PresentationGenerator:
    def __init__(self, theme: str = "dark_pro", 
                 slide_width: float = Inches(13.333), 
                 slide_height: float = Inches(7.5)):
        
        self.prs = Presentation()
        self.prs.slide_width = slide_width
        self.prs.slide_height = slide_height
        
        # Инициализация модулей
        self.theme = ColorThemes.get_theme(theme)
        self.layout_engine = SmartLayoutEngine(slide_width, slide_height)
        self.text_module = AdvancedTextModule()
        self.graphics_builder = GraphicsBuilder()
        self.media_manager = MediaManager()
        self.style_system = StyleSystem(self.theme)
        
        self.slides = []
        self.current_slide = None
        self.elements_registry = {}
        
        print(f"🚀 Инициализирован генератор презентаций (тема: {theme})")
    
    def create_slide(self, title: str = "") -> 'Slide':
        slide_layout = self.prs.slide_layouts[6]  # Пустой layout
        slide_obj = Slide(self, slide_layout, title)
        self.slides.append(slide_obj)
        self.current_slide = slide_obj
        
        # Создаем фон и заголовок
        slide_obj._create_background()
        if title:
            slide_obj._create_title()
            
        return slide_obj
    
    def add_text(self, text: str, x: float, y: float, 
                 width: float, height: float, **kwargs) -> ContentElement:
        if not self.current_slide:
            raise ValueError("Сначала создайте слайд с create_slide()")
        
        element = self.text_module.create_text_element(
            text=text, x=x, y=y, width=width, height=height, **kwargs
        )
        return self.current_slide.add_element(element)
    
    def add_shape(self, shape_type: str, x: float, y: float,
                  width: float, height: float, **kwargs) -> ContentElement:
        if not self.current_slide:
            raise ValueError("Сначала создайте слайд с create_slide()")
        
        element = self.graphics_builder.create_shape(
            shape_type=shape_type, x=x, y=y, width=width, height=height, **kwargs
        )
        return self.current_slide.add_element(element)
    
    def add_container(self, x: float, y: float, width: float, height: float, 
                      **kwargs) -> ContentElement:
        if not self.current_slide:
            raise ValueError("Сначала создайте слайд с create_slide()")
        
        element = ContentElement(
            id=f"container_{len(self.elements_registry)}",
            type=ContentType.CONTAINER,
            content={},
            x=x, y=y, width=width, height=height
        )
        
        if kwargs:
            element = self.style_system.apply_style(element, kwargs)
            
        return self.current_slide.add_element(element)
    
    def add_image(self, image_path: str, x: float, y: float, 
                  width: float, height: float, **kwargs):
        """Добавляет изображение по полному пути"""
        if not self.current_slide:
            raise ValueError("Сначала создайте слайд с create_slide()")
        
        element = self.media_manager.create_image_element(
            image_path=image_path, x=x, y=y, width=width, height=height, **kwargs
        )
        return self.current_slide.add_element(element)
    
    def add_image_by_name(self, image_name: str, x: float, y: float, 
                         width: float, height: float, **kwargs):
        """Добавляет изображение по имени из папки Картинки"""
        if not self.current_slide:
            raise ValueError("Сначала создайте слайд с create_slide()")
        
        element = self.media_manager.create_image_by_name(
            image_name=image_name, x=x, y=y, width=width, height=height, **kwargs
        )
        return self.current_slide.add_element(element)
    
    def save(self, filename: str) -> str:
        self.prs.save(filename)
        print(f"✅ Презентация сохранена: {filename}")
        return filename

class Slide:
    def __init__(self, generator: PresentationGenerator, layout, title: str):
        self.generator = generator
        self.slide = generator.prs.slides.add_slide(layout)
        self.title = title
        self.elements = []
    
    def _create_background(self):
        """Создает фон слайда"""
        try:
            background = self.slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = self.generator.theme["background"]
        except:
            # Альтернативный способ создания фона
            bg_shape = self.slide.shapes.add_shape(
                PPTX_MSO_SHAPE.RECTANGLE, 
                0, 0, 
                self.generator.prs.slide_width, 
                self.generator.prs.slide_height
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = self.generator.theme["background"]
            bg_shape.line.fill.background()
            # Перемещаем фон на задний план
            bg_shape._element.getparent().remove(bg_shape._element)
            self.slide.shapes._spTree.insert(0, bg_shape._element)
    
    def _create_title(self):
        """Создает заголовок слайда"""
        if not self.title:
            return
            
        try:
            title_box = self.slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), 
                self.generator.prs.slide_width - Inches(1), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.clear()
            p = title_frame.paragraphs[0]
            run = p.add_run()
            run.text = self.title
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = self.generator.theme["text"]
        except Exception as e:
            print(f"⚠️  Ошибка создания заголовка: {e}")
    
    def add_element(self, element: ContentElement) -> ContentElement:
        self.elements.append(element)
        self.generator.elements_registry[element.id] = element
        return element
    
    def render(self):
        for element in self.elements:
            self._render_element(element)
        return self
    
    def _render_element(self, element: ContentElement):
        bounds = self.generator.layout_engine.calculate_bounds(element)
        if not bounds:
            return
            
        x, y, width, height = bounds
        
        if element.type == ContentType.TEXT:
            self._render_text_element(element, x, y, width, height)
        elif element.type == ContentType.SHAPE:
            self._render_shape_element(element, x, y, width, height)
        elif element.type == ContentType.CONTAINER:
            self._render_container_element(element, x, y, width, height)
        elif element.type == ContentType.IMAGE:
            self._render_image_element(element, x, y, width, height)
        
        self.generator.layout_engine.reserve_area(x, y, width, height)
    
    def _render_text_element(self, element: ContentElement, x: float, y: float, 
                           width: float, height: float):
        shape = self.slide.shapes.add_shape(PPTX_MSO_SHAPE.RECTANGLE, x, y, width, height)
        self._apply_element_styles(shape, element)
        
        text_frame = shape.text_frame
        text_frame.clear()
        
        font_size = self.generator.text_module.calculate_font_size(element, width, height)
        
        text_style = element.style.text_style or TextStyle()
        text_frame.vertical_anchor = text_style.vertical_align
        p = text_frame.paragraphs[0]
        p.alignment = text_style.align
        
        run = p.add_run()
        run.text = element.content if isinstance(element.content, str) else str(element.content)
        run.font.size = Pt(font_size)
        run.font.bold = text_style.bold
        run.font.italic = text_style.italic
        if text_style.font_color:
            run.font.color.rgb = text_style.font_color
    
    def _render_shape_element(self, element: ContentElement, x: float, y: float,
                            width: float, height: float):
        shape_type = self.generator.graphics_builder._get_shape_type(
            element.content.get('shape_type', 'rectangle') if isinstance(element.content, dict) else 'rectangle'
        )
        shape = self.slide.shapes.add_shape(shape_type, x, y, width, height)
        self._apply_element_styles(shape, element)
    
    def _render_container_element(self, element: ContentElement, x: float, y: float,
                                width: float, height: float):
        shape = self.slide.shapes.add_shape(PPTX_MSO_SHAPE.RECTANGLE, x, y, width, height)
        self._apply_element_styles(shape, element)
    
    def _render_image_element(self, element: ContentElement, x: float, y: float,
                            width: float, height: float):
        """Рендерит изображение"""
        try:
            self.slide.shapes.add_picture(element.content, x, y, width, height)
        except Exception as e:
            print(f"⚠️  Ошибка загрузки изображения {element.content}: {e}")
            # Fallback - placeholder
            placeholder = self.slide.shapes.add_shape(PPTX_MSO_SHAPE.RECTANGLE, x, y, width, height)
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = self.generator.theme["secondary"]
    
    def _apply_element_styles(self, shape, element: ContentElement):
        """Безопасное применение стилей к элементу"""
        try:
            # Фон
            if element.style.background_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = element.style.background_color
            else:
                shape.fill.background()
            
            # Граница
            if element.style.border and element.style.border.color:
                shape.line.color.rgb = element.style.border.color
                
                # Безопасное преобразование ширины границы
                if element.style.border.width:
                    try:
                        # Преобразуем точки в EMU (1 point = 12700 EMU)
                        width_value = element.style.border.width
                        if hasattr(width_value, 'pt'):  # Если это объект Pt
                            width_emu = int(width_value.pt * 12700)
                        else:  # Если это число
                            width_emu = int(width_value * 12700)
                        
                        # Проверяем диапазон (0-20116800 EMU)
                        width_emu = max(0, min(20116800, width_emu))
                        shape.line.width = width_emu
                        
                    except (ValueError, TypeError) as e:
                        print(f"⚠️  Ошибка установки ширины границы: {e}")
                        shape.line.width = 12700  # Значение по умолчанию: 1 точка
            else:
                shape.line.fill.background()
                
        except Exception as e:
            print(f"⚠️  Ошибка применения стилей: {e}")
            # Устанавливаем базовые стили в случае ошибки
            shape.fill.background()
            shape.line.fill.background()

def create_presentation(theme: str = "dark_pro", **kwargs) -> PresentationGenerator:
    return PresentationGenerator(theme=theme, **kwargs)

# Экспорт
__all__ = [
    'PresentationGenerator', 'create_presentation', 'Slide',
    'ContentElement', 'ContentType', 'ElementStyle', 'TextStyle', 
    'BorderStyle', 'SizeConstraints', 'LayoutStrategy',
    'Inches', 'Pt', 'RGBColor', 'PP_ALIGN', 'MSO_SHAPE'
]