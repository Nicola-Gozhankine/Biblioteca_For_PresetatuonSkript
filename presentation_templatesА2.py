# presentation_templates_final.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from dataclasses import dataclass
from typing import List, Dict, Optional
import math

@dataclass
class ColorScheme:
    """Улучшенная цветовая схема с поддержкой темных тем"""
    name: str
    primary: RGBColor
    secondary: RGBColor
    accent: RGBColor
    background: RGBColor
    text: RGBColor
    footnote: RGBColor = None
    is_dark: bool = False
    
    def __post_init__(self):
        """Автоматическая настройка после создания"""
        if self.footnote is None:
            if self.is_dark:
                self.footnote = RGBColor(180, 180, 180)
            else:
                self.footnote = RGBColor(100, 100, 100)
    
    def get_text_color(self):
        """Автоматически выбирает контрастный цвет текста"""
        return RGBColor(255, 255, 255) if self.is_dark else self.text

class ColorThemes:
    """Расширенные цветовые темы с темными вариантами"""
    
    # Сохраняем РАБОЧИЕ темы из 2А.txt
    BLUE_TECH = ColorScheme(
        name="blue_tech",
        primary=RGBColor(41, 128, 185),
        secondary=RGBColor(52, 152, 219),
        accent=RGBColor(230, 126, 34),
        background=RGBColor(255, 255, 255),
        text=RGBColor(33, 33, 33)
    )
    
    DARK_PRO = ColorScheme(
        name="dark_pro",
        primary=RGBColor(44, 62, 80),
        secondary=RGBColor(52, 73, 94),
        accent=RGBColor(231, 76, 60),
        background=RGBColor(25, 25, 35),
        text=RGBColor(255, 255, 255),
        is_dark=True
    )
    
    GREEN_CORPORATE = ColorScheme(
        name="green_corporate",
        primary=RGBColor(39, 174, 96),
        secondary=RGBColor(46, 204, 113),
        accent=RGBColor(142, 68, 173),
        background=RGBColor(255, 255, 255),
        text=RGBColor(33, 33, 33)
    )
    
    # ДОБАВЛЯЕМ новые темные темы из 2Г.txt
    DARK_GREEN = ColorScheme(
        name="dark_green", 
        primary=RGBColor(39, 174, 96),
        secondary=RGBColor(46, 204, 113),
        accent=RGBColor(142, 68, 173),
        background=RGBColor(20, 40, 30),
        text=RGBColor(255, 255, 255),
        is_dark=True
    )
    
    DARK_PURPLE = ColorScheme(
        name="dark_purple",
        primary=RGBColor(155, 89, 182),
        secondary=RGBColor(142, 68, 173),
        accent=RGBColor(230, 126, 34),
        background=RGBColor(35, 25, 45),
        text=RGBColor(255, 255, 255),
        is_dark=True
    )
    
    @classmethod
    def get_theme(cls, theme_name: str) -> ColorScheme:
        """Безопасное получение темы с fallback"""
        themes = {
            'blue_tech': cls.BLUE_TECH,
            'dark_pro': cls.DARK_PRO,
            'green_corporate': cls.GREEN_CORPORATE,
            'dark_green': cls.DARK_GREEN,
            'dark_purple': cls.DARK_PURPLE,
        }
        return themes.get(theme_name, cls.BLUE_TECH)

class IconLibrary:
    """УЛЬТРА-БЕЗОПАСНАЯ библиотека иконок - ТОЛЬКО ТО, ЧТО ЕСТЬ В ВАШЕМ СПИСКЕ"""
    
    # ТОЛЬКО КОНСТАНТЫ ИЗ ВАШЕГО ДИАГНОСТИЧЕСКОГО СПИСКА!
    _SAFE_SHAPES = {
        # Базовые формы из 2А.txt (точно есть)
        'bullet': MSO_SHAPE.OVAL,
        'circle': MSO_SHAPE.OVAL,
        'square': MSO_SHAPE.RECTANGLE,
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rect': MSO_SHAPE.ROUNDED_RECTANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        
        # Стрелки из 2Г.txt (точно есть)
        'arrow': MSO_SHAPE.RIGHT_ARROW,
        'left_arrow': MSO_SHAPE.LEFT_ARROW,
        'up_arrow': MSO_SHAPE.UP_ARROW,
        'down_arrow': MSO_SHAPE.DOWN_ARROW,
        
        # Дополнительные формы (точно есть)
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'pentagon': MSO_SHAPE.PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5_POINT,
        'cloud': MSO_SHAPE.CLOUD,
        'heart': MSO_SHAPE.HEART,
        'lightning': MSO_SHAPE.LIGHTNING_BOLT,
        
        # УБИРАЕМ ВСЕ ПРОБЛЕМНЫЕ КОНСТАНТЫ!
        # НЕТ: 'check', 'home', 'info', 'warning' и т.д.
    }
    
    @staticmethod
    def get_safe_shape(icon_type: str) -> MSO_SHAPE:
        """Безопасное получение формы"""
        return IconLibrary._SAFE_SHAPES.get(icon_type, MSO_SHAPE.OVAL)
    
    @staticmethod
    def create_icon(slide, x, y, icon_type: str, color: RGBColor, size: str = 'medium'):
        """Создание иконки"""
        sizes = {
            'small': (Pt(20), Pt(20)),
            'medium': (Pt(30), Pt(30)),
            'large': (Pt(40), Pt(40))
        }
        width, height = sizes.get(size, (Pt(30), Pt(30)))
        
        shape_type = IconLibrary.get_safe_shape(icon_type)
        shape = slide.shapes.add_shape(shape_type, x, y, width, height)
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)
        
        return shape
    
    @staticmethod
    def create_arrow_connector(slide, start_x, start_y, end_x, end_y, color: RGBColor, width: Pt = Pt(2)):
        """Создание соединительной стрелки"""
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
        )
        connector.line.color.rgb = color
        connector.line.width = width
        
        connector.line.begin_arrowhead_width = 2
        connector.line.begin_arrowhead_length = 2
        connector.line.begin_arrowhead_style = 1
        connector.line.end_arrowhead_width = 2
        connector.line.end_arrowhead_length = 2
        connector.line.end_arrowhead_style = 1
        
        return connector

class SmartShape:
    """Умная фигура с авто-масштабированием"""
    
    def __init__(self, slide, shape_type, x, y, width, height, color_scheme: ColorScheme):
        self.slide = slide
        self.shape = slide.shapes.add_shape(shape_type, x, y, width, height)
        self.color_scheme = color_scheme
        self.text_frame = self.shape.text_frame
        self._setup_shape()
    
    def _setup_shape(self):
        """Базовая настройка фигуры"""
        self.shape.fill.solid()
        self.shape.fill.fore_color.rgb = self.color_scheme.primary
        self.shape.line.color.rgb = self.color_scheme.secondary
        self.shape.line.width = Pt(2)
    
    def set_text(self, text: str, font_size: int = None, bold: bool = False, color: RGBColor = None):
        """Установка текста с авто-масштабированием"""
        self.text_frame.clear()
        self.text_frame.vertical_anchor = 1
        p = self.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        if not font_size:
            font_size = self._calculate_font_size(text)
        
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or self.color_scheme.get_text_color()
    
    def _calculate_font_size(self, text: str) -> int:
        """Автоматический расчет размера шрифта"""
        base_size = 16
        max_length = 50
        
        if len(text) > max_length:
            return max(8, base_size - (len(text) - max_length) // 5)
        return base_size

class PresentationTemplates:
    """Основной класс библиотеки шаблонов"""
    
    def __init__(self, theme: str = 'blue_tech'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.color_scheme = ColorThemes.get_theme(theme)
        self.current_slide = None
    
    # ===== ВСЕ МЕТОДЫ ИЗ 2А.txt С УЛУЧШЕНИЯМИ ИЗ 2Г.txt =====
    
    def create_title_slide(self, title: str, subtitle: str = None, 
                          author: str = None, organization: str = None,
                          logos: List[str] = None, footnote: str = None):
        """Создание титульного слайда"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        
        self._add_background(slide)
        
        if logos:
            self._add_logos(slide, logos)
        
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), 
            self.prs.slide_width - Inches(2), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = self.color_scheme.get_text_color()
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4),
                self.prs.slide_width - Inches(2), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = subtitle
            run.font.size = Pt(28)
            run.font.color.rgb = self.color_scheme.secondary
        
        footer_y = self.prs.slide_height - Inches(1.5)
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(1), footer_y, Inches(6), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.clear()
            p = author_frame.paragraphs[0]
            run = p.add_run()
            run.text = f"Автор: {author}"
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.get_text_color()
        
        if organization:
            org_box = slide.shapes.add_textbox(
                self.prs.slide_width - Inches(7), footer_y, Inches(6), Inches(0.5)
            )
            org_frame = org_box.text_frame
            org_frame.clear()
            p = org_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = organization
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.get_text_color()
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def create_content_slide(self, title: str, sections: List[str], 
                           layout: str = 'grid', columns: int = 2,
                           icon_type: str = 'bullet', footnote: str = None):
        """Создание слайда содержания"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        if layout == 'grid':
            self._create_grid_layout(slide, sections, columns, icon_type)
        elif layout == 'list':
            self._create_list_layout(slide, sections, icon_type)
        elif layout == 'arrow_list':
            self._create_arrow_list_layout(slide, sections)
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def _create_grid_layout(self, slide, sections: List[str], columns: int, icon_type: str):
        """Сеточное расположение разделов"""
        num_sections = len(sections)
        rows = (num_sections + columns - 1) // columns
        
        section_width = (self.prs.slide_width - Inches(2)) / columns
        section_height = Inches(4.5) / rows if rows > 0 else Inches(4.5)
        
        for i, section in enumerate(sections):
            row = i // columns
            col = i % columns
            
            x = Inches(1) + col * section_width
            y = Inches(2) + row * section_height
            
            box = SmartShape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, section_width - Inches(0.5), section_height - Inches(0.3),
                self.color_scheme
            )
            box.set_text(section)
    
    def _create_list_layout(self, slide, sections: List[str], icon_type: str):
        """Списочное расположение разделов"""
        for i, section in enumerate(sections):
            y_pos = Inches(1.8) + i * Inches(0.8)
            
            IconLibrary.create_icon(
                slide, Inches(1), y_pos, icon_type, self.color_scheme.accent
            )
            
            section_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos, 
                self.prs.slide_width - Inches(2.6), Inches(0.6)
            )
            section_frame = section_box.text_frame
            section_frame.clear()
            p = section_frame.paragraphs[0]
            run = p.add_run()
            run.text = section
            run.font.size = Pt(18)
            run.font.color.rgb = self.color_scheme.get_text_color()
    
    def _create_arrow_list_layout(self, slide, sections: List[str]):
        """Список со стрелками"""
        for i, section in enumerate(sections):
            y_pos = Inches(1.8) + i * Inches(0.8)
            
            IconLibrary.create_icon(
                slide, Inches(1), y_pos, 'arrow', self.color_scheme.accent
            )
            
            section_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos, 
                self.prs.slide_width - Inches(2.6), Inches(0.6)
            )
            section_frame = section_box.text_frame
            section_frame.clear()
            p = section_frame.paragraphs[0]
            run = p.add_run()
            run.text = section
            run.font.size = Pt(18)
            run.font.color.rgb = self.color_scheme.get_text_color()
            
            if i < len(sections) - 1:
                IconLibrary.create_arrow_connector(
                    slide,
                    Inches(1.2), y_pos + Inches(0.8),
                    Inches(1.2), y_pos + Inches(1.2),
                    self.color_scheme.secondary,
                    Pt(1.5)
                )
    
    def create_info_slide(self, title: str, content_points: List[str], 
                         image_path: str = None, icon_type: str = 'bullet',
                         footnote: str = None, subheading: str = None):
        """Создание слайда с текстом и изображением"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        if subheading:
            self._add_subheading(slide, subheading)
        
        content_width = self.prs.slide_width * 0.5 - Inches(1)
        content_x = Inches(1)
        content_start_y = Inches(2.5) if subheading else Inches(2)
        
        for i, point in enumerate(content_points):
            y_pos = content_start_y + i * Inches(0.8)
            
            IconLibrary.create_icon(
                slide, content_x, y_pos, icon_type, self.color_scheme.accent
            )
            
            text_box = slide.shapes.add_textbox(
                content_x + Inches(0.6), y_pos - Inches(0.1),
                content_width - Inches(0.6), Inches(0.6)
            )
            text_frame = text_box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(16)
            run.font.color.rgb = self.color_scheme.get_text_color()
        
        if image_path:
            img_width = self.prs.slide_width * 0.5 - Inches(1)
            img_x = self.prs.slide_width * 0.5
            img_height = self.prs.slide_height - Inches(2)
            
            try:
                slide.shapes.add_picture(
                    image_path, img_x, Inches(1), img_width, img_height
                )
            except:
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, img_x, Inches(1), img_width, img_height
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = self.color_scheme.secondary
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def create_infographic_slide(self, title: str, infographic_type: str, 
                               data: Dict, footnote: str = None, **kwargs):
        """Создание слайда инфографики"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        if infographic_type == 'pyramid':
            self._create_pyramid(slide, data, **kwargs)
        elif infographic_type == 'circles':
            self._create_connected_circles(slide, data, **kwargs)
        elif infographic_type == 'flowchart':
            self._create_flowchart(slide, data, **kwargs)
        elif infographic_type == 'arrow_flow':
            self._create_arrow_flowchart(slide, data, **kwargs)
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def _create_connected_circles(self, slide, data: Dict, radius: float = 1.5):
        """Создание связанных кругов"""
        center_x = self.prs.slide_width / 2
        center_y = self.prs.slide_height / 2 + Inches(0.5)
        
        items = data.get('items', [])
        num_items = len(items)
        
        if num_items == 0:
            return
            
        circles = []
        
        for i, item in enumerate(items):
            angle = 2 * math.pi * i / num_items
            x = center_x + Inches(radius) * math.cos(angle) - Inches(0.8)
            y = center_y + Inches(radius) * math.sin(angle) - Inches(0.8)
            
            circle = SmartShape(
                slide, MSO_SHAPE.OVAL,
                x, y, Inches(1.6), Inches(1.6), self.color_scheme
            )
            text = item.get('text', item) if isinstance(item, dict) else item
            circle.set_text(text)
            center_circle_x = x + Inches(0.8)
            center_circle_y = y + Inches(0.8)
            circles.append((circle, center_circle_x, center_circle_y))
        
        for i in range(len(circles)):
            for j in range(i + 1, len(circles)):
                circle1, x1, y1 = circles[i]
                circle2, x2, y2 = circles[j]
                
                IconLibrary.create_arrow_connector(
                    slide, x1, y1, x2, y2, self.color_scheme.secondary
                )
    
    def _create_pyramid(self, slide, data: Dict):
        """Создание пирамиды"""
        items = data.get('items', [])
        num_levels = len(items)
        
        if num_levels == 0:
            return
            
        pyramid_width = Inches(8)
        level_height = Inches(1)
        start_x = (self.prs.slide_width - pyramid_width) / 2
        start_y = Inches(2)
        
        for i, item in enumerate(items):
            level_width = pyramid_width * (i + 1) / num_levels
            level_x = start_x + (pyramid_width - level_width) / 2
            level_y = start_y + (num_levels - i - 1) * level_height
            
            level = SmartShape(
                slide, MSO_SHAPE.RECTANGLE,
                level_x, level_y, level_width, level_height - Inches(0.1),
                self.color_scheme
            )
            text = item.get('text', item) if isinstance(item, dict) else item
            level.set_text(text)
    
    def _create_flowchart(self, slide, data: Dict):
        """Создание блок-схемы"""
        items = data.get('items', [])
        
        for i, item in enumerate(items):
            x = Inches(2)
            y = Inches(1.8) + i * Inches(1.2)
            
            box = SmartShape(
                slide, MSO_SHAPE.RECTANGLE,
                x, y, Inches(4), Inches(0.8),
                self.color_scheme
            )
            text = item.get('text', item) if isinstance(item, dict) else item
            box.set_text(text)
            
            if i < len(items) - 1:
                IconLibrary.create_arrow_connector(
                    slide,
                    x + Inches(2), y + Inches(0.8),
                    x + Inches(2), y + Inches(1.2),
                    self.color_scheme.secondary
                )
    
    def _create_arrow_flowchart(self, slide, data: Dict):
        """Блок-схема с большими стрелками"""
        items = data.get('items', [])
        
        for i, item in enumerate(items):
            x = Inches(3)
            y = Inches(1.8) + i * Inches(1.5)
            
            box = SmartShape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, Inches(4), Inches(0.8),
                self.color_scheme
            )
            text = item.get('text', item) if isinstance(item, dict) else item
            box.set_text(text)
            
            if i < len(items) - 1:
                arrow_y = y + Inches(1)
                IconLibrary.create_icon(
                    slide, x + Inches(1.5), arrow_y, 'arrow', 
                    self.color_scheme.accent, 'large'
                )
    
    def create_sources_slide(self, title: str, sources: List[str], 
                           image_path: str = None, footnote: str = None):
        """Создание слайда с источниками"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        for i, source in enumerate(sources):
            y_pos = Inches(1.5) + i * Inches(0.7)
            
            number = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), y_pos, Inches(0.4), Inches(0.4)
            )
            number.fill.solid()
            number.fill.fore_color.rgb = self.color_scheme.accent
            
            num_text = slide.shapes.add_textbox(Inches(1), y_pos, Inches(0.4), Inches(0.4))
            num_frame = num_text.text_frame
            num_frame.vertical_anchor = 1
            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = self.color_scheme.get_text_color()
            
            source_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos, Inches(6), Inches(0.6)
            )
            source_frame = source_box.text_frame
            source_frame.clear()
            p = source_frame.paragraphs[0]
            run = p.add_run()
            run.text = source
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.get_text_color()
        
        if image_path:
            img_width = self.prs.slide_width * 0.4
            img_x = self.prs.slide_width * 0.55
            img_height = self.prs.slide_height - Inches(2)
            
            try:
                slide.shapes.add_picture(
                    image_path, img_x, Inches(1), img_width, img_height
                )
            except:
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, img_x, Inches(1), img_width, img_height
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = self.color_scheme.secondary
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    # ===== ВСПОМОГАТЕЛЬНЫЕ МЕТОДЫ =====
    def _add_background(self, slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.color_scheme.background
    
    def _add_title(self, slide, title: str):
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), 
            self.prs.slide_width - Inches(1), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = self.color_scheme.get_text_color()
    
    def _add_subheading(self, slide, subheading: str):
        subheading_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2),
            self.prs.slide_width - Inches(1), Inches(0.6)
        )
        subheading_frame = subheading_box.text_frame
        subheading_frame.clear()
        p = subheading_frame.paragraphs[0]
        run = p.add_run()
        run.text = subheading
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = self.color_scheme.secondary
    
    def _add_footnote(self, slide, footnote: str):
        footnote_box = slide.shapes.add_textbox(
            Inches(0.5), self.prs.slide_height - Inches(0.8),
            self.prs.slide_width - Inches(1), Inches(0.5)
        )
        footnote_frame = footnote_box.text_frame
        footnote_frame.clear()
        p = footnote_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = footnote
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = self.color_scheme.footnote
    
    def _add_logos(self, slide, logos: List[str]):
        logo_size = Inches(1.2)
        spacing = (self.prs.slide_width - len(logos) * logo_size) / (len(logos) + 1)
        
        for i, logo_path in enumerate(logos):
            x = spacing + i * (logo_size + spacing)
            try:
                slide.shapes.add_picture(logo_path, x, Inches(0.3), logo_size, logo_size)
            except:
                logo = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(0.3), logo_size, logo_size
                )
                logo.fill.solid()
                logo.fill.fore_color.rgb = self.color_scheme.primary
    
    def save(self, filename: str):
        self.prs.save(filename)
        print(f"✅ Презентация сохранена: {filename}")

# ===== ПРОСТОЙ ИНТЕРФЕЙС =====
def create_presentation(slides_config: List[Dict], theme: str = 'blue_tech', 
                       output_file: str = 'presentation.pptx'):
    template = PresentationTemplates(theme)
    
    for slide_config in slides_config:
        config = slide_config.copy()
        slide_type = config.pop('type')
        
        if slide_type == 'title':
            template.create_title_slide(**config)
        elif slide_type == 'content':
            template.create_content_slide(**config)
        elif slide_type == 'info':
            template.create_info_slide(**config)
        elif slide_type == 'infographic':
            template.create_infographic_slide(**config)
        elif slide_type == 'sources':
            template.create_sources_slide(**config)
        else:
            print(f"⚠️  Неизвестный тип слайда: {slide_type}")
    
    template.save(output_file)
    return output_file