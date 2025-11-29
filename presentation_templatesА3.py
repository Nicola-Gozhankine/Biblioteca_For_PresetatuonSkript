# presentation_library.py
"""
УНИФИЦИРОВАННАЯ БИБЛИОТЕКА ДЛЯ СОЗДАНИЯ ПРЕЗЕНТАЦИЙ
Версия 1.0 - Единый файл со всеми модулями
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Union, Tuple
from enum import Enum
import math

# ===== БАЗОВЫЕ ТИПЫ ДАННЫХ =====
class ContentType(Enum):
    TEXT = "text"
    IMAGE = "image"
    SHAPE = "shape"
    CONTAINER = "container"
    INFOGRAPHIC = "infographic"

class LayoutStrategy(Enum):
    MANUAL = "manual"
    GRID = "grid"
    VERTICAL_STACK = "vstack"
    HORIZONTAL_STACK = "hstack"

@dataclass
class SizeConstraints:
    min_width: Optional[float] = None
    max_width: Optional[float] = None
    min_height: Optional[float] = None
    max_height: Optional[float] = None
    aspect_ratio: Optional[float] = None
    grow_priority: int = 1

@dataclass
class BorderStyle:
    color: Optional[RGBColor] = None
    width: float = Pt(1)
    radius: Optional[float] = None

@dataclass
class TextStyle:
    font_size: Optional[float] = None
    font_color: Optional[RGBColor] = None
    bold: bool = False
    italic: bool = False
    align: PP_ALIGN = PP_ALIGN.LEFT
    vertical_align: int = 1

@dataclass
class ElementStyle:
    background_color: Optional[RGBColor] = None
    border: Optional[BorderStyle] = None
    text_style: Optional[TextStyle] = None
    padding: float = Inches(0.1)
    margin: float = Inches(0.05)

@dataclass
class ContentElement:
    id: str
    type: ContentType
    content: Union[str, Dict, List]
    style: ElementStyle = field(default_factory=ElementStyle)
    constraints: SizeConstraints = field(default_factory=SizeConstraints)
    layout_strategy: LayoutStrategy = LayoutStrategy.MANUAL
    parent: Optional['ContentElement'] = None
    children: List['ContentElement'] = field(default_factory=list)
    x: Optional[float] = None
    y: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None
    metadata: Dict = field(default_factory=dict)
    
    def add_child(self, child: 'ContentElement'):
        child.parent = self
        self.children.append(child)
        return self

# ===== СИСТЕМА ТЕМ =====
class ColorThemes:
    """Цветовые темы презентаций"""
    
    THEMES = {
        "default": {
            "primary": RGBColor(41, 128, 185),
            "secondary": RGBColor(52, 152, 219),
            "accent": RGBColor(230, 126, 34),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(33, 33, 33),
            "success": RGBColor(39, 174, 96),
            "warning": RGBColor(241, 196, 15),
            "danger": RGBColor(231, 76, 60),
            "footnote": RGBColor(100, 100, 100)
        },
        "dark": {
            "primary": RGBColor(44, 62, 80),
            "secondary": RGBColor(52, 73, 94),
            "accent": RGBColor(231, 76, 60),
            "background": RGBColor(25, 25, 35),
            "text": RGBColor(255, 255, 255),
            "success": RGBColor(39, 174, 96),
            "warning": RGBColor(241, 196, 15),
            "danger": RGBColor(231, 76, 60),
            "footnote": RGBColor(180, 180, 180)
        },
        "blue_tech": {
            "primary": RGBColor(41, 128, 185),
            "secondary": RGBColor(52, 152, 219),
            "accent": RGBColor(230, 126, 34),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(33, 33, 33),
            "success": RGBColor(39, 174, 96),
            "warning": RGBColor(241, 196, 15),
            "danger": RGBColor(231, 76, 60),
            "footnote": RGBColor(100, 100, 100)
        },
        "green_corporate": {
            "primary": RGBColor(39, 174, 96),
            "secondary": RGBColor(46, 204, 113),
            "accent": RGBColor(142, 68, 173),
            "background": RGBColor(255, 255, 255),
            "text": RGBColor(33, 33, 33),
            "success": RGBColor(39, 174, 96),
            "warning": RGBColor(241, 196, 15),
            "danger": RGBColor(231, 76, 60),
            "footnote": RGBColor(100, 100, 100)
        }
    }
    
    @classmethod
    def get_theme(cls, theme_name: str) -> Dict:
        return cls.THEMES.get(theme_name, cls.THEMES["default"])
    
    @classmethod
    def get_available_themes(cls) -> List[str]:
        return list(cls.THEMES.keys())

# ===== МОДУЛЬ ТЕКСТА =====
class TextModule:
    """Модуль работы с текстом и авто-масштабированием"""
    
    def __init__(self):
        # НАСТРОЙКИ РАЗМЕРОВ ШРИФТА ПО УМОЛЧАНИЮ
        self.FONT_CONFIG = {
            "title": {"base": 22, "min": 18, "max": 32},
            "subtitle": {"base": 20, "min": 18, "max": 24},
            "main": {"base": 16, "min": 14, "max": 18},  # ОСНОВНОЙ ТЕКСТ 16-18
            "caption": {"base": 12, "min": 10, "max": 14},
            "footnote": {"base": 10, "min": 8, "max": 12}
        }
    
    def calculate_font_size(self, element: ContentElement, available_width: float, available_height: float) -> float:
        """АВТО-МАСШТАБИРОВАНИЕ ШРИФТА С ПРАВИЛЬНЫМИ МИНИМУМАМИ"""
        if not isinstance(element.content, str):
            return self.FONT_CONFIG["main"]["base"]
        
        text = element.content
        text_type = self._detect_text_type(element, text)
        config = self.FONT_CONFIG[text_type]
        
        base_size = config["base"]
        min_size = config["min"]
        max_size = config["max"]
        
        # УЧЕТ ДЛИНЫ ТЕКСТА
        chars_per_inch = 7  # Оптимальная плотность символов
        max_chars_per_line = available_width * chars_per_inch
        
        if len(text) > max_chars_per_line and max_chars_per_line > 0:
            lines_needed = len(text) / max_chars_per_line
            # Мягкое уменьшение для длинного текста
            reduction = min(4, (lines_needed - 1) * 1.5)
            base_size = max(min_size, base_size - reduction)
        
        # УЧЕТ ВЫСОТЫ
        line_height = base_size * 0.02  # Примерная высота строки в дюймах
        available_lines = available_height / line_height if line_height > 0 else 1
        
        if available_lines < 1:
            # Если не помещается по высоте, уменьшаем размер
            height_based_size = available_height / 0.02
            base_size = max(min_size, min(base_size, height_based_size))
        
        # ФИНАЛЬНЫЕ ОГРАНИЧЕНИЯ
        final_size = max(min_size, min(max_size, base_size))
        return final_size
    
    def _detect_text_type(self, element: ContentElement, text: str) -> str:
        """Определяет тип текста для правильного подбора размера"""
        
        # Явные указания в стилях
        if element.style.text_style:
            if element.style.text_style.font_size:
                if element.style.text_style.font_size >= 20:
                    return "title"
                elif element.style.text_style.font_size <= 10:
                    return "footnote"
            
            if element.style.text_style.bold and len(text) < 50:
                return "title"
        
        # Эвристики на основе содержания
        text_lower = text.lower()
        
        # Очень короткий текст
        if len(text) < 15:
            return "caption"
        
        # Текст с маркерами или цифрами
        if any(marker in text for marker in ['•', '-', '—', '·', '1.', '2.', '3.']):
            return "subtitle"
        
        # Длинный текст
        if len(text) > 100:
            return "main"
        
        # Заголовки обычно короткие
        if len(text) < 50 and (text.isupper() or text[0].isupper() and text[-1] in '!?:'):
            return "title"
        
        return "main"

# ===== МОДУЛЬ КОМПОНОВКИ =====
class LayoutModule:
    """Модуль управления размещением элементов"""
    
    def __init__(self, slide_width: float, slide_height: float):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.occupied_areas = []
        self.safe_margin = Inches(0.5)
    
    def calculate_bounds(self, element: ContentElement, parent_bounds: Optional[Tuple] = None) -> Optional[Tuple]:
        """Вычисляет границы элемента"""
        if element.layout_strategy == LayoutStrategy.MANUAL:
            return self._calculate_manual_bounds(element, parent_bounds)
        # TODO: Добавить другие стратегии
        return self._calculate_manual_bounds(element, parent_bounds)
    
    def _calculate_manual_bounds(self, element: ContentElement, parent_bounds: Optional[Tuple]) -> Tuple:
        """Ручное размещение"""
        if parent_bounds:
            parent_x, parent_y, parent_w, parent_h = parent_bounds
            x = parent_x + (element.x or 0)
            y = parent_y + (element.y or 0)
            width = min(element.width or parent_w, parent_w)
            height = min(element.height or parent_h, parent_h)
        else:
            x = element.x or self.safe_margin
            y = element.y or self.safe_margin
            width = element.width or (self.slide_width - 2 * self.safe_margin)
            height = element.height or (self.slide_height - 2 * self.safe_margin)
        
        return (x, y, width, height)
    
    def reserve_area(self, x: float, y: float, width: float, height: float):
        """Резервирует область"""
        self.occupied_areas.append((x, y, width, height))
    
    def calculate_child_bounds(self, parent_bounds: Tuple, padding: float) -> Tuple:
        """Вычисляет границы для дочернего элемента"""
        parent_x, parent_y, parent_w, parent_h = parent_bounds
        return (
            parent_x + padding,
            parent_y + padding,
            parent_w - 2 * padding,
            parent_h - 2 * padding
        )

# ===== ОСНОВНОЙ API =====
class PresentationBuilder:
    """
    ОСНОВНОЙ КЛАСС ДЛЯ СОЗДАНИЯ ПРЕЗЕНТАЦИЙ
    Простой API для быстрого старта
    """
    
    def __init__(self, theme: str = "default"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = ColorThemes.get_theme(theme)
        self.text_module = TextModule()
        self.layout_module = LayoutModule(self.prs.slide_width, self.prs.slide_height)
        self.current_slide = None
        self.elements_registry = {}
    
    def create_slide(self, title: str = "") -> 'SlideBuilder':
        """Создает новый слайд"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        
        # Фон
        self._create_background(slide)
        
        # Заголовок
        if title:
            self._create_title(slide, title)
        
        return SlideBuilder(self, slide)
    
    def _create_background(self, slide):
        """Создает фон слайда"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme["background"]
    
    def _create_title(self, slide, title: str):
        """Создает заголовок слайда"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), 
            self.prs.slide_width - Inches(1), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = self.theme["text"]
    
    def save(self, filename: str):
        """Сохраняет презентацию"""
        self.prs.save(filename)
        print(f"✅ Презентация сохранена: {filename}")

class SlideBuilder:
    """Строитель слайдов"""
    
    def __init__(self, builder: PresentationBuilder, slide):
        self.builder = builder
        self.slide = slide
        self.elements = []
        self.layout_module = LayoutModule(builder.prs.slide_width, builder.prs.slide_height)
    
    def add_element(self, element: ContentElement) -> 'SlideBuilder':
        """Добавляет элемент на слайд"""
        self.elements.append(element)
        self.builder.elements_registry[element.id] = element
        return self
    
    def render(self) -> 'SlideBuilder':
        """Рендерит все элементы на слайд"""
        for element in self.elements:
            self._render_element(element)
        return self
    
    def _render_element(self, element: ContentElement, parent_bounds: Optional[Tuple] = None):
        """Рендерит элемент и его детей"""
        bounds = self.layout_module.calculate_bounds(element, parent_bounds)
        if not bounds:
            return
        
        x, y, width, height = bounds
        
        # Рендерим элемент в зависимости от типа
        if element.type == ContentType.TEXT:
            self._render_text_element(element, x, y, width, height)
        elif element.type == ContentType.CONTAINER:
            self._render_container_element(element, x, y, width, height)
        elif element.type == ContentType.SHAPE:
            self._render_shape_element(element, x, y, width, height)
        
        # Резервируем область
        self.layout_module.reserve_area(x, y, width, height)
        
        # Рекурсивно рендерим детей
        for child in element.children:
            child_bounds = self.layout_module.calculate_child_bounds(
                (x, y, width, height), element.style.padding
            )
            self._render_element(child, child_bounds)
    
    def _render_text_element(self, element: ContentElement, x: float, y: float, width: float, height: float):
        """Рендерит текстовый элемент"""
        # Создаем форму
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        
        # Применяем стили
        self._apply_element_styles(shape, element)
        
        # Настраиваем текст
        text_frame = shape.text_frame
        text_frame.clear()
        
        # АВТО-МАСШТАБИРОВАНИЕ через текстовый модуль
        font_size = self.builder.text_module.calculate_font_size(element, width, height)
        
        # Текстовые стили
        text_style = element.style.text_style or TextStyle()
        
        text_frame.vertical_anchor = text_style.vertical_align
        p = text_frame.paragraphs[0]
        p.alignment = text_style.align
        
        run = p.add_run()
        run.text = element.content if isinstance(element.content, str) else str(element.content)
        run.font.size = Pt(font_size)
        run.font.bold = text_style.bold
        run.font.italic = text_style.italic
        run.font.color.rgb = text_style.font_color or self.builder.theme["text"]
    
    def _render_container_element(self, element: ContentElement, x: float, y: float, width: float, height: float):
        """Рендерит контейнер"""
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        self._apply_element_styles(shape, element)
    
    def _render_shape_element(self, element: ContentElement, x: float, y: float, width: float, height: float):
        """Рендерит геометрическую фигуру"""
        shape_type = self._get_shape_type(element.content.get('shape_type', 'rectangle') if isinstance(element.content, dict) else 'rectangle')
        shape = self.slide.shapes.add_shape(shape_type, x, y, width, height)
        self._apply_element_styles(shape, element)
    
    def _apply_element_styles(self, shape, element: ContentElement):
        """Применяет стили к форме"""
        # Фон
        if element.style.background_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = element.style.background_color
        else:
            shape.fill.background()
        
        # Граница
        if element.style.border and element.style.border.color:
            shape.line.color.rgb = element.style.border.color
            shape.line.width = element.style.border.width
        else:
            shape.line.fill.background()
    
    def _get_shape_type(self, shape_name: str) -> MSO_SHAPE:
        """Возвращает тип фигуры по имени"""
        shapes = {
            'rectangle': MSO_SHAPE.RECTANGLE,
            'rounded_rect': MSO_SHAPE.ROUNDED_RECTANGLE,
            'oval': MSO_SHAPE.OVAL,
            'circle': MSO_SHAPE.OVAL,
            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
            'diamond': MSO_SHAPE.DIAMOND,
            'star': MSO_SHAPE.STAR_5_POINT,
        }
        return shapes.get(shape_name, MSO_SHAPE.RECTANGLE)

# ===== ПРОСТОЙ ИНТЕРФЕЙС =====
def create_presentation(theme: str = "default") -> PresentationBuilder:
    """Создает новую презентацию"""
    return PresentationBuilder(theme=theme)

def create_text_element(text: str, x: float = None, y: float = None, 
                       width: float = None, height: float = None, 
                       element_id: str = None) -> ContentElement:
    """Быстрое создание текстового элемента"""
    element_id = element_id or f"text_{id(text)}"
    return ContentElement(
        id=element_id,
        type=ContentType.TEXT,
        content=text,
        x=x,
        y=y,
        width=width,
        height=height
    )

def create_container_element(x: float = None, y: float = None,
                           width: float = None, height: float = None,
                           element_id: str = None) -> ContentElement:
    """Быстрое создание контейнера"""
    element_id = element_id or f"container_{id(element_id)}"
    return ContentElement(
        id=element_id,
        type=ContentType.CONTAINER,
        content={},
        x=x,
        y=y,
        width=width,
        height=height
    )

def create_shape_element(shape_type: str, x: float = None, y: float = None,
                        width: float = None, height: float = None,
                        element_id: str = None) -> ContentElement:
    """Быстрое создание геометрической фигуры"""
    element_id = element_id or f"shape_{shape_type}"
    return ContentElement(
        id=element_id,
        type=ContentType.SHAPE,
        content={'shape_type': shape_type},
        x=x,
        y=y,
        width=width,
        height=height
    )

# ===== ТЕСТИРОВАНИЕ =====
if __name__ == "__main__":
    """Тестирование библиотеки"""
    print("🧪 ТЕСТИРОВАНИЕ БИБЛИОТЕКИ...")
    
    try:
        # 1. Создание презентации
        presentation = create_presentation(theme="blue_tech")
        print("✅ Создание презентации - УСПЕХ")
        
        # 2. Создание слайда
        slide = presentation.create_slide("Тестовый слайд")
        print("✅ Создание слайда - УСПЕХ")
        
        # 3. Создание элементов
        title = create_text_element(
            "Главный заголовок",
            x=Inches(1), y=Inches(1), 
            width=Inches(5), height=Inches(0.8),
            element_id="main_title"
        )
        title.style.background_color = presentation.theme["primary"]
        title.style.text_style = TextStyle(
            font_color=RGBColor(255, 255, 255),
            bold=True,
            align=PP_ALIGN.CENTER
        )
        
        content = create_text_element(
            "Это тестовый контент для проверки работы авто-масштабирования шрифта. "
            "Текст должен быть читаемым и хорошо смотреться в рамке.",
            x=Inches(1), y=Inches(2), 
            width=Inches(5), height=Inches(1.5),
            element_id="content"
        )
        content.style.border = BorderStyle(
            color=presentation.theme["accent"],
            width=Pt(2)
        )
        content.style.padding = Inches(0.2)
        
        # Контейнер с вложенным элементом
        container = create_container_element(
            x=Inches(7), y=Inches(1),
            width=Inches(4), height=Inches(3),
            element_id="main_container"
        )
        container.style.border = BorderStyle(
            color=presentation.theme["secondary"],
            width=Pt(3)
        )
        container.style.padding = Inches(0.3)
        
        inner_text = create_text_element(
            "Вложенный текст",
            element_id="inner_text"
        )
        inner_text.style.background_color = presentation.theme["success"]
        inner_text.style.text_style = TextStyle(
            font_color=RGBColor(255, 255, 255),
            align=PP_ALIGN.CENTER
        )
        
        container.add_child(inner_text)
        
        # 4. Добавление элементов на слайд
        slide.add_element(title).add_element(content).add_element(container)
        print("✅ Добавление элементов - УСПЕХ")
        
        # 5. Рендеринг
        slide.render()
        print("✅ Рендеринг слайда - УСПЕХ")
        
        # 6. Сохранение
        presentation.save("test_presentation_library.pptx")
        print("✅ Сохранение презентации - УСПЕХ")
        
        print("\n🎉 БИБЛИОТЕКА РАБОТАЕТ КОРРЕКТНО!")
        print("📁 test_presentation_library.pptx - создан для проверки")
        
    except Exception as e:
        print(f"❌ ОШИБКА: {e}")
        import traceback
        traceback.print_exc()