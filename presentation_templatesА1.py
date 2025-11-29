# presentation_templates.py
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from dataclasses import dataclass
from typing import List, Dict, Tuple, Optional
import math

@dataclass
class ColorScheme:
    """Цветовая схема презентации"""
    name: str
    primary: RGBColor
    secondary: RGBColor
    accent: RGBColor
    background: RGBColor
    text: RGBColor
    success: RGBColor = None
    warning: RGBColor = None
    danger: RGBColor = None

class ColorThemes:
    """Предустановленные цветовые темы"""
    
    BLUE_TECH = ColorScheme(
        name="blue_tech",
        primary=RGBColor(41, 128, 185),
        secondary=RGBColor(52, 152, 219),
        accent=RGBColor(230, 126, 34),
        background=RGBColor(255, 255, 255),
        text=RGBColor(33, 33, 33)
    )
    
    DARK_PRO = ColorScheme(
        name="dark_pro",
        primary=RGBColor(44, 62, 80),
        secondary=RGBColor(52, 73, 94),
        accent=RGBColor(231, 76, 60),
        background=RGBColor(25, 25, 35),
        text=RGBColor(255, 255, 255)
    )
    
    GREEN_CORPORATE = ColorScheme(
        name="green_corporate",
        primary=RGBColor(39, 174, 96),
        secondary=RGBColor(46, 204, 113),
        accent=RGBColor(142, 68, 173),
        background=RGBColor(255, 255, 255),
        text=RGBColor(33, 33, 33)
    )
    
    @classmethod
    def get_theme(cls, theme_name: str) -> ColorScheme:
        themes = {
            'blue_tech': cls.BLUE_TECH,
            'dark_pro': cls.DARK_PRO,
            'green_corporate': cls.GREEN_CORPORATE
        }
        return themes.get(theme_name, cls.BLUE_TECH)

class IconLibrary:
    """Библиотека стандартных иконок и фигур"""
    
    @staticmethod
    def create_icon(shape, icon_type: str, color: RGBColor):
        """Создание стандартных иконок"""
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)
        return shape

class SmartShape:
    """Умная фигура с авто-масштабированием текста"""
    
    def __init__(self, slide, shape_type, x, y, width, height, color_scheme: ColorScheme):
        self.slide = slide
        self.shape = slide.shapes.add_shape(shape_type, x, y, width, height)
        self.color_scheme = color_scheme
        self.text_frame = self.shape.text_frame
        self._setup_shape()
    
    def _setup_shape(self):
        """Базовая настройка фигуры"""
        self.shape.fill.solid()
        self.shape.fill.fore_color.rgb = self.color_scheme.primary
        self.shape.line.color.rgb = self.color_scheme.secondary
        self.shape.line.width = Pt(2)
    
    def set_text(self, text: str, font_size: int = None, bold: bool = False):
        """Установка текста с авто-масштабированием"""
        self.text_frame.clear()
        self.text_frame.vertical_anchor = 1  # Middle
        p = self.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        # Авто-подбор размера шрифта
        if not font_size:
            font_size = self._calculate_font_size(text)
        
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = self.color_scheme.text
    
    def _calculate_font_size(self, text: str) -> int:
        """Автоматический расчет размера шрифта на основе длины текста"""
        base_size = 16
        max_length = 50
        
        if len(text) > max_length:
            return max(8, base_size - (len(text) - max_length) // 5)
        return base_size

class PresentationTemplates:
    """Основной класс библиотеки шаблонов"""
    
    def __init__(self, theme: str = 'blue_tech'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.color_scheme = ColorThemes.get_theme(theme)
        self.current_slide = None
    
    # ===== ШАБЛОН 1: ТИТУЛЬНЫЙ СЛАЙД =====
    def create_title_slide(self, title: str, subtitle: str = None, 
                          author: str = None, organization: str = None,
                          logos: List[str] = None):
        """Создание титульного слайда"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        
        # Фон
        self._add_background(slide)
        
        # Эмблемы сверху (если есть)
        if logos:
            self._add_logos(slide, logos)
        
        # Основной заголовок
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), 
            self.prs.slide_width - Inches(2), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = self.color_scheme.text
        
        # Подзаголовок
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4),
                self.prs.slide_width - Inches(2), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = subtitle
            run.font.size = Pt(28)
            run.font.color.rgb = self.color_scheme.secondary
        
        # Подписи внизу
        footer_y = self.prs.slide_height - Inches(1.5)
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(1), footer_y, Inches(6), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_frame.clear()
            p = author_frame.paragraphs[0]
            run = p.add_run()
            run.text = f"Автор: {author}"
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.text
        
        if organization:
            org_box = slide.shapes.add_textbox(
                self.prs.slide_width - Inches(7), footer_y, Inches(6), Inches(0.5)
            )
            org_frame = org_box.text_frame
            org_frame.clear()
            p = org_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = organization
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.text
        
        return slide
    
    # ===== ШАБЛОН 2: СОДЕРЖАНИЕ =====
    def create_content_slide(self, title: str, sections: List[str], 
                           layout: str = 'grid', columns: int = 2):
        """Создание слайда содержания"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        # Заголовок
        self._add_title(slide, title)
        
        # Размещение разделов в зависимости от layout
        if layout == 'grid':
            self._create_grid_layout(slide, sections, columns)
        elif layout == 'list':
            self._create_list_layout(slide, sections)
        
        return slide
    
    def _create_grid_layout(self, slide, sections: List[str], columns: int):
        """Сеточное расположение разделов"""
        num_sections = len(sections)
        rows = (num_sections + columns - 1) // columns
        
        section_width = (self.prs.slide_width - Inches(2)) / columns
        section_height = Inches(4.5) / rows if rows > 0 else Inches(4.5)
        
        for i, section in enumerate(sections):
            row = i // columns
            col = i % columns
            
            x = Inches(1) + col * section_width
            y = Inches(2) + row * section_height
            
            # Создаем рамку для раздела
            box = SmartShape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, section_width - Inches(0.5), section_height - Inches(0.3),
                self.color_scheme
            )
            box.set_text(section)
    
    def _create_list_layout(self, slide, sections: List[str]):
        """Списочное расположение разделов"""
        for i, section in enumerate(sections):
            y_pos = Inches(1.8) + i * Inches(0.8)
            
            # Иконка
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), y_pos, Inches(0.4), Inches(0.4)
            )
            IconLibrary.create_icon(icon, 'bullet', self.color_scheme.accent)
            
            # Текст раздела
            section_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos, 
                self.prs.slide_width - Inches(2.6), Inches(0.6)
            )
            section_frame = section_box.text_frame
            section_frame.clear()
            p = section_frame.paragraphs[0]
            run = p.add_run()
            run.text = section
            run.font.size = Pt(18)
            run.font.color.rgb = self.color_scheme.text
    
    # ===== ШАБЛОН 3: ОБЩАЯ ИНФОРМАЦИЯ =====
    def create_info_slide(self, title: str, content_points: List[str], 
                         image_path: str = None, icon_type: str = 'bullet'):
        """Создание слайда с текстом и изображением"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        # Заголовок
        self._add_title(slide, title)
        
        # Область для контента (50% слева)
        content_width = self.prs.slide_width * 0.5 - Inches(1)
        content_x = Inches(1)
        
        # Добавляем пункты с иконками
        for i, point in enumerate(content_points):
            y_pos = Inches(2) + i * Inches(0.8)
            
            # Иконка
            icon = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, content_x, y_pos, Inches(0.4), Inches(0.4)
            )
            IconLibrary.create_icon(icon, icon_type, self.color_scheme.accent)
            
            # Текст
            text_box = slide.shapes.add_textbox(
                content_x + Inches(0.6), y_pos - Inches(0.1),
                content_width - Inches(0.6), Inches(0.6)
            )
            text_frame = text_box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(16)
            run.font.color.rgb = self.color_scheme.text
        
        # Изображение справа (50%)
        if image_path:
            img_width = self.prs.slide_width * 0.5 - Inches(1)
            img_x = self.prs.slide_width * 0.5
            img_height = self.prs.slide_height - Inches(2)
            
            try:
                slide.shapes.add_picture(
                    image_path, img_x, Inches(1), img_width, img_height
                )
            except:
                # Запасной вариант если изображение не найдено
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, img_x, Inches(1), img_width, img_height
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = self.color_scheme.secondary
        
        return slide
    
    # ===== ШАБЛОН 4: ИНФОГРАФИКА =====
    def create_infographic_slide(self, title: str, infographic_type: str, 
                               data: Dict, **kwargs):
        """Создание слайда инфографики"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        # Выбор типа инфографики
        if infographic_type == 'pyramid':
            self._create_pyramid(slide, data, **kwargs)
        elif infographic_type == 'circles':
            self._create_connected_circles(slide, data, **kwargs)
        elif infographic_type == 'flowchart':
            self._create_flowchart(slide, data, **kwargs)
        
        return slide
    
    def _create_connected_circles(self, slide, data: Dict, radius: float = 1.5):
        """Создание связанных кругов"""
        center_x = self.prs.slide_width / 2
        center_y = self.prs.slide_height / 2 + Inches(0.5)
        
        items = data.get('items', [])
        num_items = len(items)
        
        if num_items == 0:
            return
            
        circles = []
        
        for i, item in enumerate(items):
            angle = 2 * math.pi * i / num_items
            x = center_x + Inches(radius) * math.cos(angle) - Inches(0.8)
            y = center_y + Inches(radius) * math.sin(angle) - Inches(0.8)
            
            # Создаем круг
            circle = SmartShape(
                slide, MSO_SHAPE.OVAL,
                x, y, Inches(1.6), Inches(1.6), self.color_scheme
            )
            circle.set_text(item.get('text', item) if isinstance(item, dict) else item)
            center_circle_x = x + Inches(0.8)
            center_circle_y = y + Inches(0.8)
            circles.append((circle, center_circle_x, center_circle_y))
        
        # Соединяем круги линиями
        for i in range(len(circles)):
            for j in range(i + 1, len(circles)):
                circle1, x1, y1 = circles[i]
                circle2, x2, y2 = circles[j]
                
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2
                )
                connector.line.color.rgb = self.color_scheme.secondary
                connector.line.width = Pt(1.5)
    
    def _create_pyramid(self, slide, data: Dict):
        """Создание пирамиды"""
        items = data.get('items', [])
        num_levels = len(items)
        
        if num_levels == 0:
            return
            
        pyramid_width = Inches(8)
        level_height = Inches(1)
        start_x = (self.prs.slide_width - pyramid_width) / 2
        start_y = Inches(2)
        
        for i, item in enumerate(items):
            level_width = pyramid_width * (i + 1) / num_levels
            level_x = start_x + (pyramid_width - level_width) / 2
            level_y = start_y + (num_levels - i - 1) * level_height
            
            level = SmartShape(
                slide, MSO_SHAPE.RECTANGLE,
                level_x, level_y, level_width, level_height - Inches(0.1),
                self.color_scheme
            )
            level.set_text(item.get('text', item) if isinstance(item, dict) else item)
    
    def _create_flowchart(self, slide, data: Dict):
        """Создание блок-схемы"""
        items = data.get('items', [])
        
        for i, item in enumerate(items):
            x = Inches(2)
            y = Inches(1.8) + i * Inches(1.2)
            
            box = SmartShape(
                slide, MSO_SHAPE.RECTANGLE,
                x, y, Inches(4), Inches(0.8),
                self.color_scheme
            )
            box.set_text(item.get('text', item) if isinstance(item, dict) else item)
            
            # Соединительные линии (кроме последнего элемента)
            if i < len(items) - 1:
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    x + Inches(2), y + Inches(0.8),
                    x + Inches(2), y + Inches(1.2)
                )
                connector.line.color.rgb = self.color_scheme.secondary
                connector.line.width = Pt(2)
    
    # ===== ШАБЛОН 5: ИСТОЧНИКИ =====
    def create_sources_slide(self, title: str, sources: List[str], 
                           image_path: str = None):
        """Создание слайда с источниками"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.current_slide = slide
        self._add_background(slide)
        
        self._add_title(slide, title)
        
        # Список источников слева
        for i, source in enumerate(sources):
            y_pos = Inches(1.5) + i * Inches(0.7)
            
            # Нумерация
            number = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1), y_pos, Inches(0.4), Inches(0.4)
            )
            number.fill.solid()
            number.fill.fore_color.rgb = self.color_scheme.accent
            
            # Номер
            num_text = slide.shapes.add_textbox(Inches(1), y_pos, Inches(0.4), Inches(0.4))
            num_frame = num_text.text_frame
            num_frame.vertical_anchor = 1
            p = num_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = self.color_scheme.text
            
            # Источник
            source_box = slide.shapes.add_textbox(
                Inches(1.6), y_pos, Inches(6), Inches(0.6)
            )
            source_frame = source_box.text_frame
            source_frame.clear()
            p = source_frame.paragraphs[0]
            run = p.add_run()
            run.text = source
            run.font.size = Pt(14)
            run.font.color.rgb = self.color_scheme.text
        
        # Изображение справа
        if image_path:
            img_width = self.prs.slide_width * 0.4
            img_x = self.prs.slide_width * 0.55
            img_height = self.prs.slide_height - Inches(2)
            
            try:
                slide.shapes.add_picture(
                    image_path, img_x, Inches(1), img_width, img_height
                )
            except:
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, img_x, Inches(1), img_width, img_height
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = self.color_scheme.secondary
        
        return slide
    
    # ===== ВСПОМОГАТЕЛЬНЫЕ МЕТОДЫ =====
    def _add_background(self, slide):
        """Добавление фона слайда"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.color_scheme.background
    
    def _add_title(self, slide, title: str):
        """Добавление заголовка слайда"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), 
            self.prs.slide_width - Inches(1), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        p = title_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = self.color_scheme.text
    
    def _add_logos(self, slide, logos: List[str]):
        """Добавление логотипов"""
        logo_size = Inches(1.2)
        spacing = (self.prs.slide_width - len(logos) * logo_size) / (len(logos) + 1)
        
        for i, logo_path in enumerate(logos):
            x = spacing + i * (logo_size + spacing)
            try:
                slide.shapes.add_picture(logo_path, x, Inches(0.3), logo_size, logo_size)
            except:
                # Запасной вариант
                logo = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(0.3), logo_size, logo_size
                )
                logo.fill.solid()
                logo.fill.fore_color.rgb = self.color_scheme.primary
    
    def save(self, filename: str):
        """Сохранение презентации"""
        self.prs.save(filename)
        print(f"✅ Презентация сохранена: {filename}")

# ===== ИСПРАВЛЕННЫЙ ПРОСТОЙ ИНТЕРФЕЙС =====
def create_presentation(slides_config: List[Dict], theme: str = 'blue_tech', 
                       output_file: str = 'presentation.pptx'):
    """Простой интерфейс для создания презентации"""
    
    template = PresentationTemplates(theme)
    
    for slide_config in slides_config:
        # Создаем копию конфигурации без ключа 'type'
        config = slide_config.copy()
        slide_type = config.pop('type')  # Убираем 'type' из аргументов
        
        if slide_type == 'title':
            template.create_title_slide(**config)
        elif slide_type == 'content':
            template.create_content_slide(**config)
        elif slide_type == 'info':
            template.create_info_slide(**config)
        elif slide_type == 'infographic':
            template.create_infographic_slide(**config)
        elif slide_type == 'sources':
            template.create_sources_slide(**config)
        else:
            print(f"⚠️  Неизвестный тип слайда: {slide_type}")
    
    template.save(output_file)
    return output_file